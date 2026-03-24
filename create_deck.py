#!/usr/bin/env python3
"""
NextStepUni Year 2 Proposal Deck — PwC Empowering Futures Programme
Creates a clean, modern 12-slide PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────
ACCENT      = RGBColor(0xCC, 0x78, 0x5C)   # warm terracotta
DARK        = RGBColor(0x1A, 0x1A, 0x1A)   # titles
BODY        = RGBColor(0x33, 0x33, 0x33)   # body text
SECONDARY   = RGBColor(0x66, 0x66, 0x66)   # subtitles / notes
LIGHT_BG    = RGBColor(0xF7, 0xF5, 0xF3)   # light warm background for cards
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
VERY_LIGHT  = RGBColor(0xF0, 0xF0, 0xF0)   # table alt row

FONT_NAME = "Calibri"

# ── Dimensions ──────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]  # blank

# ── Helpers ─────────────────────────────────────────────────────

def add_accent_bar(slide, top=False):
    """Thin 2pt accent line across the top of the slide."""
    if top:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            SLIDE_W, Pt(4)
        )
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            SLIDE_W, Pt(4)
        )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=16,
                bold=False, color=BODY, alignment=PP_ALIGN.LEFT,
                font_name=FONT_NAME, line_spacing=None, space_after=None):
    """Add a simple textbox and return (textbox_shape, text_frame)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    if space_after is not None:
        p.space_after = Pt(space_after)
    return txBox, tf


def add_paragraph(tf, text, font_size=16, bold=False, color=BODY,
                  alignment=PP_ALIGN.LEFT, space_before=0, space_after=0,
                  line_spacing=None, font_name=FONT_NAME):
    """Append a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return p


def add_title(slide, text, left=Inches(0.9), top=Inches(0.6),
              width=Inches(11), size=32):
    """Standard slide title."""
    _, tf = add_textbox(slide, left, top, width, Inches(0.7),
                        text, font_size=size, bold=True, color=DARK)
    return tf


def add_subtitle(slide, text, left=Inches(0.9), top=Inches(1.2),
                 width=Inches(11), size=16):
    """Subtitle line under the title."""
    add_textbox(slide, left, top, width, Inches(0.5),
                text, font_size=size, color=SECONDARY)


def add_bottom_note(slide, text, top=Inches(6.6)):
    """Small note at the bottom of the slide."""
    add_textbox(slide, Inches(0.9), top, Inches(11), Inches(0.5),
                text, font_size=12, color=SECONDARY)


def add_rect_card(slide, left, top, width, height, fill_color=LIGHT_BG):
    """A rounded-corner-ish card background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Gentle rounding
    shape.adjustments[0] = 0.05
    return shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)

# Large accent bar at top for title slide
bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(0), Inches(0), SLIDE_W, Inches(0.25))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

# Vertical accent stripe on left
stripe = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.9), Inches(2.2), Inches(0.06), Inches(3.0))
stripe.fill.solid()
stripe.fill.fore_color.rgb = ACCENT
stripe.line.fill.background()

add_textbox(sl, Inches(1.3), Inches(2.2), Inches(8), Inches(1.0),
            "NextStepUni", font_size=48, bold=True, color=DARK)

add_textbox(sl, Inches(1.3), Inches(3.1), Inches(8), Inches(0.6),
            "Year 1 Update & Year 2 Evolution", font_size=24, color=BODY)

add_textbox(sl, Inches(1.3), Inches(3.8), Inches(8), Inches(0.5),
            "PwC Empowering Futures Programme", font_size=16, color=SECONDARY)

add_textbox(sl, Inches(1.3), Inches(4.7), Inches(4), Inches(0.4),
            "February 2026", font_size=14, color=SECONDARY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Year 1 at a Glance
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "Year 1: Building the Foundation")

# Three stat cards
stats = [
    ("6", "NEIC Schools"),
    ("~500", "Students"),
    ("100%", "Schools Actively Engaged"),
]
card_w = Inches(3.2)
card_h = Inches(2.2)
gap = Inches(0.5)
total_w = card_w * 3 + gap * 2
start_x = (SLIDE_W - total_w) / 2

for i, (num, label) in enumerate(stats):
    x = start_x + i * (card_w + gap)
    y = Inches(1.9)
    add_rect_card(sl, x, y, card_w, card_h)
    add_textbox(sl, x, y + Inches(0.35), card_w, Inches(0.9),
                num, font_size=44, bold=True, color=ACCENT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(sl, x, y + Inches(1.3), card_w, Inches(0.6),
                label, font_size=16, color=BODY,
                alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.2), Inches(4.7), Inches(10.5), Inches(1.2),
            "In-person workshops delivered across all six schools. Students onboarded onto the "
            "online study programme. Guidance counsellor feedback collected throughout.",
            font_size=15, color=SECONDARY, alignment=PP_ALIGN.CENTER,
            line_spacing=24)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — What We've Learned
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "What Year 1 Has Taught Us")

learnings = [
    "Students respond best when they can interact with the material, not just watch it",
    "Personalisation matters — one-size-fits-all content doesn't land with every student",
    "Guidance counsellors want tools students can use between sessions",
    "Real impact data is essential for tracking what's working",
]

for i, item in enumerate(learnings):
    y = Inches(2.0) + Inches(i * 1.1)
    # Small accent dot
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(1.2), y + Inches(0.12), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    add_textbox(sl, Inches(1.65), y, Inches(9.5), Inches(0.6),
                item, font_size=17, color=BODY, line_spacing=24)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Direct Feedback from the Schools
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "Direct Feedback from the Schools")
add_subtitle(sl, "From guidance counsellor feedback forms")

quotes = [
    "\u201cStudents need something interactive — not just videos they\u2019ll watch once and forget\u201d",
    "\u201cA tool they can use on their own, that meets them where they are\u201d",
    "\u201cSomething that connects study skills to their actual subjects and goals\u201d",
]

card_w = Inches(10.5)
card_h = Inches(1.0)
start_y = Inches(2.2)

for i, q in enumerate(quotes):
    y = start_y + Inches(i * 1.3)
    add_rect_card(sl, Inches(1.2), y, card_w, card_h)
    # Quote mark accent
    add_textbox(sl, Inches(1.5), y + Inches(0.05), Inches(0.5), Inches(0.7),
                "\u201c", font_size=36, bold=True, color=ACCENT)
    add_textbox(sl, Inches(2.0), y + Inches(0.2), Inches(9.2), Inches(0.7),
                q.strip("\u201c\u201d"), font_size=16, color=BODY)

add_bottom_note(sl, "Collected via structured feedback forms from participating guidance counsellors")


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Year 1 → Year 2: The Evolution
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "Year 1 \u2192 Year 2: The Evolution")

col_w = Inches(4.8)
col_h = Inches(4.0)
col1_x = Inches(1.0)
col2_x = Inches(7.0)
col_y = Inches(1.8)

# Left card
add_rect_card(sl, col1_x, col_y, col_w, col_h)
add_textbox(sl, col1_x + Inches(0.4), col_y + Inches(0.25), col_w - Inches(0.8), Inches(0.5),
            "Year 1: Foundation", font_size=18, bold=True, color=ACCENT)

left_items = [
    "In-person workshops",
    "Online video course",
    "End-of-section quizzes",
    "General content",
    "Completion tracking",
]
_, tf_l = add_textbox(sl, col1_x + Inches(0.4), col_y + Inches(0.9),
                      col_w - Inches(0.8), Inches(3.0),
                      left_items[0], font_size=15, color=BODY)
tf_l.paragraphs[0].space_after = Pt(10)
for item in left_items[1:]:
    add_paragraph(tf_l, item, font_size=15, color=BODY, space_after=10)

# Right card
add_rect_card(sl, col2_x, col_y, col_w, col_h)
add_textbox(sl, col2_x + Inches(0.4), col_y + Inches(0.25), col_w - Inches(0.8), Inches(0.5),
            "Year 2: Full Platform", font_size=18, bold=True, color=ACCENT)

right_items = [
    "In-person workshops (stays)",
    "Interactive learning platform",
    "Built-in study tools & simulators",
    "Personalised to each student",
    "Full engagement & impact data",
]
_, tf_r = add_textbox(sl, col2_x + Inches(0.4), col_y + Inches(0.9),
                      col_w - Inches(0.8), Inches(3.0),
                      right_items[0], font_size=15, color=BODY)
tf_r.paragraphs[0].space_after = Pt(10)
for item in right_items[1:]:
    add_paragraph(tf_r, item, font_size=15, color=BODY, space_after=10)

# Arrow between columns
arrow = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                            Inches(5.95), Inches(3.4), Inches(0.9), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT
arrow.line.fill.background()

add_bottom_note(sl, "Everything from Year 1 carries forward. Year 2 adds the layer that turns knowledge into action.")


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — The Platform
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "The NextStepUni Learning Platform")

features = [
    ("52 Interactive Modules",
     "Covering mindset, learning science, subject-specific\nstrategies, and exam preparation"),
    ("Personalised Onboarding",
     "Students enter their subjects, current grades,\nand target grades"),
    ("Built-in Study Tools",
     "Exam planners, breathing exercises,\nstudy schedulers, grade calculators"),
    ("AI Study Assistant",
     "Powered by Google Gemini for\npersonalised study support"),
]

card_w = Inches(5.2)
card_h = Inches(1.8)
gap_x = Inches(0.6)
gap_y = Inches(0.4)
start_x = Inches(1.0)
start_y = Inches(1.8)

for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_rect_card(sl, x, y, card_w, card_h)
    # Accent bar on left side of card
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              x, y, Inches(0.06), card_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_textbox(sl, x + Inches(0.35), y + Inches(0.2), card_w - Inches(0.6), Inches(0.5),
                title, font_size=17, bold=True, color=DARK)
    add_textbox(sl, x + Inches(0.35), y + Inches(0.7), card_w - Inches(0.6), Inches(1.0),
                desc, font_size=14, color=SECONDARY, line_spacing=20)

add_bottom_note(sl, "Live demo available during this meeting", top=Inches(6.5))


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Built for These Students
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "Built for These Students")

differentiators = [
    ("Research-Backed", "Every module built on established educational psychology"),
    ("Socioeconomically Conscious", "No assumptions about resources, tutors, or quiet study spaces"),
    ("Irish Context", "Leaving Cert, CAO points, mock exams — not generic study advice"),
    ("Guidance Counsellor Informed", "Built directly from teacher feedback"),
]

for i, (headline, desc) in enumerate(differentiators):
    y = Inches(1.9) + Inches(i * 1.2)
    # Accent line to the left
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(1.2), y, Inches(0.06), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_textbox(sl, Inches(1.6), y, Inches(9.5), Inches(0.45),
                headline, font_size=18, bold=True, color=DARK)
    add_textbox(sl, Inches(1.6), y + Inches(0.45), Inches(9.5), Inches(0.4),
                desc, font_size=15, color=SECONDARY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — The Science Behind It
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "The Science Behind It")
add_subtitle(sl, "Full detail in the supporting document")

evidence = [
    ("Up to 50%", "increase in retention", "Active learning vs passive lecture formats"),
    ("25–30%", "higher scores", "Students using spaced repetition on delayed tests"),
    ("3–5x", "higher completion", "Interactive environments vs passive video content"),
]

card_w = Inches(3.2)
card_h = Inches(2.8)
gap = Inches(0.6)
total_w = card_w * 3 + gap * 2
start_x = (SLIDE_W - total_w) / 2

for i, (stat, label, detail) in enumerate(evidence):
    x = start_x + i * (card_w + gap)
    y = Inches(2.2)
    add_rect_card(sl, x, y, card_w, card_h)
    add_textbox(sl, x, y + Inches(0.3), card_w, Inches(0.7),
                stat, font_size=36, bold=True, color=ACCENT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(sl, x, y + Inches(1.1), card_w, Inches(0.5),
                label, font_size=16, bold=True, color=DARK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(sl, x + Inches(0.3), y + Inches(1.7), card_w - Inches(0.6), Inches(0.9),
                detail, font_size=13, color=SECONDARY,
                alignment=PP_ALIGN.CENTER, line_spacing=18)

add_bottom_note(sl, "Supporting document: 'NextStepUni — Programme Overview & Evidence Base'")


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Impact Data PwC Can Report
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "Measurable Impact for PwC")
add_subtitle(sl, "What the platform tracks — and what you can report")

# Left column
col_w = Inches(5.0)
col_h = Inches(3.5)
left_x = Inches(1.0)
right_x = Inches(6.8)
col_y = Inches(2.0)

add_rect_card(sl, left_x, col_y, col_w, col_h)
add_textbox(sl, left_x + Inches(0.4), col_y + Inches(0.2), col_w - Inches(0.8), Inches(0.5),
            "Student-Level Data", font_size=17, bold=True, color=ACCENT)

left_data = [
    "Modules completed",
    "Study tools used",
    "Subject-specific engagement",
    "Individual progress over time",
]
_, tf_l = add_textbox(sl, left_x + Inches(0.4), col_y + Inches(0.8),
                      col_w - Inches(0.8), Inches(2.5),
                      "", font_size=15, color=BODY)
tf_l.paragraphs[0].text = left_data[0]
tf_l.paragraphs[0].font.size = Pt(15)
tf_l.paragraphs[0].font.color.rgb = BODY
tf_l.paragraphs[0].font.name = FONT_NAME
tf_l.paragraphs[0].space_after = Pt(8)
for item in left_data[1:]:
    p = add_paragraph(tf_l, item, font_size=15, color=BODY, space_after=8)

# Right column
add_rect_card(sl, right_x, col_y, col_w, col_h)
add_textbox(sl, right_x + Inches(0.4), col_y + Inches(0.2), col_w - Inches(0.8), Inches(0.5),
            "Programme-Level Data", font_size=17, bold=True, color=ACCENT)

right_data = [
    "Overall completion rates",
    "Total engagement hours",
    "Most-used tools and modules",
    "Before/after grade trajectories",
]
_, tf_r = add_textbox(sl, right_x + Inches(0.4), col_y + Inches(0.8),
                      col_w - Inches(0.8), Inches(2.5),
                      "", font_size=15, color=BODY)
tf_r.paragraphs[0].text = right_data[0]
tf_r.paragraphs[0].font.size = Pt(15)
tf_r.paragraphs[0].font.color.rgb = BODY
tf_r.paragraphs[0].font.name = FONT_NAME
tf_r.paragraphs[0].space_after = Pt(8)
for item in right_data[1:]:
    add_paragraph(tf_r, item, font_size=15, color=BODY, space_after=8)

add_bottom_note(sl, "First time this programme will have granular, reportable impact data")


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Year 2 Investment
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "Year 2 Investment")

# Build the table
rows, cols = 5, 3
tbl_w = Inches(10.0)
tbl_h = Inches(2.8)
tbl_left = Inches(1.5)
tbl_top = Inches(1.8)

table_shape = sl.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h)
table = table_shape.table

# Column widths
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(3.5)
table.columns[2].width = Inches(3.5)

data = [
    ["", "Year 1", "Year 2"],
    ["Schools", "6 NEIC schools", "6 NEIC schools"],
    ["Students", "~500", "~500–600"],
    ["Delivery", "Workshops + video course", "Workshops + interactive platform"],
    ["Annual Investment", "€18,000", "€30,000"],
]

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONT_NAME
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = BODY
            paragraph.alignment = PP_ALIGN.LEFT
        # Style header row
        if r == 0:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
        else:
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = VERY_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
        # Bold the row labels
        if c == 0 and r > 0:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.color.rgb = DARK
        # Accent the Year 2 column
        if c == 2 and r > 0:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
        # Last row — investment amounts — accent colour
        if r == 4 and c >= 1:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.color.rgb = ACCENT
                paragraph.font.size = Pt(16)

# Notes below table
notes_y = Inches(4.8)
_, tf_n = add_textbox(sl, Inches(1.5), notes_y, Inches(10), Inches(1.2),
                      "The \u20ac12,000 increase covers platform infrastructure, AI tools, and ongoing maintenance.",
                      font_size=14, color=BODY)
add_paragraph(tf_n,
              "The platform itself has been developed at no additional cost to PwC.",
              font_size=14, color=BODY, space_before=6)

# Accent bottom line
accent_line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.5), Inches(6.15), Inches(10), Pt(2))
accent_line.fill.solid()
accent_line.fill.fore_color.rgb = ACCENT
accent_line.line.fill.background()

add_textbox(sl, Inches(1.5), Inches(6.3), Inches(10), Inches(0.4),
            "Equivalent platform development would cost \u20ac40,000\u201360,000 to commission externally.",
            font_size=14, bold=True, color=ACCENT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — The North Star
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "Why We\u2019re Here")

statement = (
    "These students deserve the same tools and support as anyone else. "
    "Year 1 got us into the room. Year 2 gives them something that changes "
    "how they study, how they think about their future, and how they prepare "
    "for the Leaving Cert."
)

add_textbox(sl, Inches(1.5), Inches(2.5), Inches(10.3), Inches(3.0),
            statement, font_size=22, color=DARK,
            alignment=PP_ALIGN.CENTER, line_spacing=36)

# Subtle accent line below the statement
line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           Inches(5.5), Inches(5.5), Inches(2.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Next Steps
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_accent_bar(sl)
add_title(sl, "Next Steps")

steps = [
    ("1", "Agree Year 2 delivery model and investment"),
    ("2", "Begin platform onboarding for Year 2 cohort"),
    ("3", "Year 1 review meeting — June 2026"),
    ("4", "Year 2 launch — October 2026"),
]

for i, (num, text) in enumerate(steps):
    y = Inches(2.0) + Inches(i * 1.0)
    # Number circle
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(1.2), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    # Number text
    tf_c = circle.text_frame
    tf_c.paragraphs[0].text = num
    tf_c.paragraphs[0].font.size = Pt(16)
    tf_c.paragraphs[0].font.bold = True
    tf_c.paragraphs[0].font.color.rgb = WHITE
    tf_c.paragraphs[0].font.name = FONT_NAME
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_c.word_wrap = False

    add_textbox(sl, Inches(2.0), y + Inches(0.05), Inches(9), Inches(0.5),
                text, font_size=17, color=BODY)

# Divider
div = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(1.2), Inches(6.0), Inches(10.5), Pt(1.5))
div.fill.solid()
div.fill.fore_color.rgb = VERY_LIGHT
div.line.fill.background()

# Thank you + contact
add_textbox(sl, Inches(1.2), Inches(6.2), Inches(10.5), Inches(0.4),
            "Thank you", font_size=18, bold=True, color=DARK,
            alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1.2), Inches(6.6), Inches(10.5), Inches(0.4),
            "Alex Linehan  |  NextStepUni / PwC Empowering Futures",
            font_size=14, color=SECONDARY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
output_path = "/Users/alexlinehan/Downloads/NextStepUni_Year2_Proposal.pptx"
prs.save(output_path)
print(f"Saved → {output_path}")
print(f"Slides: {len(prs.slides)}")
